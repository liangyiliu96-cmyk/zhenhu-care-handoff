"""臻护平台完整演示PPT v2.0 — 基于46份项目文档，覆盖架构/功能/技术/质量/规划。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 配色 ──
PRIMARY    = RGBColor(0x0B, 0x64, 0x72)
PRIMARY_LT = RGBColor(0x14, 0x92, 0xA6)
ACCENT     = RGBColor(0xE8, 0x6A, 0x17)
SUCCESS    = RGBColor(0x2E, 0x7D, 0x32)
WARNING    = RGBColor(0xED, 0x6C, 0x02)
DANGER     = RGBColor(0xD3, 0x2F, 0x2F)
PURPLE     = RGBColor(0x7B, 0x1F, 0xA2)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG   = RGBColor(0xF5, 0xF9, 0xFA)
BORDER     = RGBColor(0xE0, 0xE0, 0xE0)
DARK_BG    = RGBColor(0x1A, 0x23, 0x2E)
HIGHLIGHT  = RGBColor(0xE8, 0xF0, 0xF2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

OUTPUT = os.path.join(os.path.dirname(__file__), '..', 'output', '臻护平台完整演示.pptx')
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)

# ── 工具函数 ──
def bg(slide, color): slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else: s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else: s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font; p.alignment = align
    return tb

def title_bar(slide, text, sub=None):
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.07), PRIMARY)
    txt(slide, Inches(0.7), Inches(0.3), Inches(11), Inches(0.55), text, 30, True, PRIMARY)
    if sub: txt(slide, Inches(0.7), Inches(0.9), Inches(11), Inches(0.35), sub, 12, False, GRAY)
    rect(slide, Inches(0.7), Inches(1.4), Inches(11.8), Inches(0.015), BORDER)

def mini_card(slide, l, t, w, h, icon, title, desc, bg_c=None, accent=None):
    c = rrect(slide, l, t, w, h, bg_c or WHITE, BORDER)
    txt(slide, l+Inches(0.18), t+Inches(0.12), w-Inches(0.36), Inches(0.38), icon, 24, False, accent or PRIMARY)
    txt(slide, l+Inches(0.18), t+Inches(0.5), w-Inches(0.36), Inches(0.28), title, 13, True, BLACK)
    txt(slide, l+Inches(0.18), t+Inches(0.8), w-Inches(0.36), h-Inches(1.0), desc, 9, False, GRAY)

def kpi_card(slide, l, t, w, h, num, label, color=PRIMARY):
    c = rrect(slide, l, t, w, h, LIGHT_BG, BORDER)
    txt(slide, l, t+Inches(0.15), w, Inches(0.45), num, 26, True, color, PP_ALIGN.CENTER)
    txt(slide, l, t+Inches(0.6), w, Inches(0.3), label, 10, False, GRAY, PP_ALIGN.CENTER)

def flow_step(slide, l, t, w, h, n, title, desc, color=PRIMARY):
    rrect(slide, l, t, w, h, color)
    txt(slide, l, t+Inches(0.08), w, Inches(0.3), str(n), 16, True, WHITE, PP_ALIGN.CENTER)
    txt(slide, l, t+Inches(0.4), w, Inches(0.25), title, 11, True, WHITE, PP_ALIGN.CENTER)
    txt(slide, l+Inches(0.08), t+Inches(0.7), w-Inches(0.16), Inches(0.6), desc, 8, False, RGBColor(0xDD,0xEE,0xF2), PP_ALIGN.CENTER)

def feature_list(slide, l, t, w, items, color=None, size=11):
    for i, (icon, title) in enumerate(items):
        y = t + Inches(i * 0.45)
        txt(slide, l, y, Inches(0.3), Inches(0.35), icon, size, False, color or PRIMARY)
        txt(slide, l+Inches(0.35), y, w-Inches(0.35), Inches(0.35), title, size, False, BLACK)

def section_card(slide, l, t, w, h, title, items, accent=None):
    bg = rrect(slide, l, t, w, h, HIGHLIGHT, BORDER)
    txt(slide, l+Inches(0.15), t+Inches(0.12), w-Inches(0.3), Inches(0.3), title, 15, True, accent or PRIMARY)
    rect(slide, l+Inches(0.15), t+Inches(0.48), w-Inches(0.3), Inches(0.015), accent or PRIMARY)
    for i, item in enumerate(items):
        y = t + Inches(0.6) + Inches(i * 0.28)
        txt(slide, l+Inches(0.25), y, w-Inches(0.4), Inches(0.25), f'• {item}', 10, False, BLACK)

# ═══════════════════════════════════════════════════════
# P1: 封面
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, PRIMARY)
rect(s, Inches(0), Inches(0), Inches(0.12), SLIDE_H, PRIMARY_LT)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.12), PRIMARY_LT)
txt(s, Inches(1.2), Inches(1.5), Inches(11), Inches(0.9), '臻 护', 64, True, WHITE)
txt(s, Inches(1.2), Inches(2.5), Inches(11), Inches(0.5), 'ZhenHu — AI 驱动的出院交接与慢病智能协同管理平台', 24, False, PRIMARY_LT)
rect(s, Inches(1.2), Inches(3.2), Inches(3.5), Inches(0.04), ACCENT)
txt(s, Inches(1.2), Inches(3.6), Inches(10), Inches(0.5), '面向科室的智能住院管理解决方案  ·  心内科试点  ·  v0.3.0', 14, False, RGBColor(0xBB,0xD5,0xDC))
txt(s, Inches(1.2), Inches(6.2), Inches(8), Inches(0.3), '基于46份设计文档  ·  163+ API端点  ·  19节点AI Agent  ·  FHIR R4标准化  ·  全角色覆盖', 10, False, GRAY)

# ═══════════════════════════════════════════════════════
# P2: 目录
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '演讲目录')
agenda = [('01', '项目背景与痛点分析'), ('02', '产品定位与解决方案'), ('03', '系统架构全景'), ('04', 'AI Agent 编排引擎'), ('05', '医生端功能详解'), ('06', '护士端功能详解'), ('07', '管理端功能详解'), ('08', 'FHIR 标准化与数据互通'), ('09', '临床智能助手体系'), ('10', '中西医协同创新'), ('11', '知识库与证据图谱'), ('12', '安全与审计体系'), ('13', '质量保障与测试'), ('14', '项目价值与预期效果'), ('15', '实施路线图'), ('16', '团队与致谢')]
for i, (num, title) in enumerate(agenda):
    col = i % 4; row = i // 4
    x = Inches(0.6) + col * Inches(3.1)
    y = Inches(2.0) + row * Inches(1.35)
    rect(s, x, y, Inches(2.8), Inches(1.1), LIGHT_BG if col in [0,2] else WHITE, BORDER)
    txt(s, x+Inches(0.15), y+Inches(0.1), Inches(2.5), Inches(0.3), num, 22, True, PRIMARY)
    txt(s, x+Inches(0.15), y+Inches(0.55), Inches(2.5), Inches(0.4), title, 12, False, BLACK)

# ═══════════════════════════════════════════════════════
# P3: 痛点分析
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '01  项目背景：住院管理的四大核心痛点')
pains = [
    ('🔴', '出院交接断裂', '出院小结与社区/下级医院信息脱节\n30%患者出院后7天内管理"真空"\n医嘱口头传递，遗漏率高达25%', DANGER),
    ('🟠', '多角色协同低效', '医护药三方依赖纸质交接班本\n信息遗漏率25%，沟通成本高\n缺少统一的临床决策平台', WARNING),
    ('🟡', 'AI能力未嵌入临床', 'NEWS2/DDx/出院标准全靠经验\n缺少结构化证据链和实时告警\n知识库无法辅助临床推理', ACCENT),
    ('🔵', '数据孤岛无法互通', 'HIS/LIS/EMR系统间数据不互通\n出院后慢病数据断裂\n无法形成完整闭环', PRIMARY),
]
for i, (icon, title, desc, color) in enumerate(pains):
    x = Inches(0.5) + i * Inches(3.2)
    rrect(s, x, Inches(1.8), Inches(2.9), Inches(3.4), WHITE, BORDER)
    txt(s, x+Inches(0.15), Inches(1.95), Inches(2.6), Inches(0.4), f'{icon}  {title}', 16, True, color)
    rect(s, x+Inches(0.15), Inches(2.45), Inches(2.0), Inches(0.02), color)
    txt(s, x+Inches(0.15), Inches(2.65), Inches(2.6), Inches(2.2), desc, 12, False, GRAY)

# 统计数据
stats = [('30%', '出院后管理真空率'), ('25%', '交接信息遗漏率'), ('14种', '覆盖病种'), ('163+', 'API端点')]
for i, (n, l) in enumerate(stats):
    kpi_card(s, Inches(1.0)+i*Inches(3.0), Inches(5.5), Inches(2.5), Inches(1.2), n, l, ACCENT)

# ═══════════════════════════════════════════════════════
# P4: 产品定位
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '02  产品定位与核心价值', '面向医护人员的临床决策支持系统 (CDSS)')

# 三端定位
roles = [
    ('👨‍⚕️', '医生端', '工作台·查房·审核·用药·出院\nAI辅助临床推理与决策\n结构化SOAP + DDx'),
    ('👩‍⚕️', '护士端', '看板·任务·监测·交班·制度\nAI优先级排序 + 体征趋势\n护理全流程闭环'),
    ('⚙️', '管理端', '知识库治理·病种模板·质控\n科主任/护士长双角色视图\n运维面板 + 证据图谱'),
]
for i, (icon, title, desc) in enumerate(roles):
    x = Inches(0.6) + i * Inches(4.2)
    rrect(s, x, Inches(1.8), Inches(3.8), Inches(2.6), WHITE, BORDER)
    txt(s, x+Inches(0.2), Inches(1.95), Inches(3.4), Inches(0.4), f'{icon}  {title}', 20, True, PRIMARY)
    txt(s, x+Inches(0.2), Inches(2.5), Inches(3.4), Inches(1.5), desc, 12, False, GRAY)

# 下方独特性
uniques = [
    ('🤖 AI Agent 编排', '19个临床节点 StateGraph\n自动模式+人工审核双模'),
    ('🏥 FHIR R4 标准', '8类资源 + Patient Compartment\nPII脱敏 + 访问审计'),
    ('🔄 版本化状态机', 'CAS乐观锁防并发冲突\n409冲突后保留用户草稿'),
    ('📊 证据链溯源', 'RAG + LLM 双重验证\nNeo4j 证据图谱可视化'),
]
for i, (icon, title) in enumerate(uniques):
    x = Inches(0.6) + i * Inches(3.2)
    mini_card(s, x, Inches(4.8), Inches(2.9), Inches(1.5), icon, title, '', WHITE)

# ═══════════════════════════════════════════════════════
# P5: 系统架构全景
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '03  系统架构全景', '五层微服务 + 六大基础设施 + 三端协同')

# 五层架构
layers = [
    ('表现层', 'React 18 + MUI v6 · Vite 6 · TanStack Query · Recharts · Zustand', PRIMARY),
    ('API 网关', 'FastAPI · CORS · Auth(header/jwt/oidc) · Idempotency · 请求ID追踪', PRIMARY_LT),
    ('核心服务', '住院协同(8000) · FHIR适配(8300) · 知识编排(8200) · 工作流引擎(8100)', ACCENT),
    ('AI 引擎', 'LangGraph Agent · DeepSeek V4 · Ollama回退 · RAG · CircuitBreaker · Harness护栏', PURPLE),
    ('数据层', 'PostgreSQL · SQLite · Neo4j证据图谱 · Redis缓存 · Milvus向量库', SUCCESS),
]
for i, (name, desc, color) in enumerate(layers):
    y = Inches(1.7) + i * Inches(0.85)
    rrect(s, Inches(0.6), y, Inches(8.5), Inches(0.7), color)
    txt(s, Inches(0.8), y+Inches(0.12), Inches(2.2), Inches(0.45), name, 16, True, WHITE)
    txt(s, Inches(3.1), y+Inches(0.12), Inches(5.5), Inches(0.45), desc, 10, False, WHITE)

# 右侧技术栈
section_card(s, Inches(9.5), Inches(1.7), Inches(3.2), Inches(5.5), '技术栈',
    ['Python 3.13 + FastAPI', 'TypeScript 5.7 + React 18', 'LangGraph 19节点编排', 'Milvus 2.4 向量检索',
     'Neo4j 证据图谱', 'Redis 7 缓存/队列', 'Docker Compose部署', 'SQLAlchemy 2.0 ORM',
     'Playwright E2E', 'Vitest + pytest'], PRIMARY_LT)

# ═══════════════════════════════════════════════════════
# P6: AI Agent 编排引擎
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '04  AI Agent 编排引擎', '19节点 StateGraph · 双模式(自动+人工审核) · 6道安全护栏')

# 19节点横向展示
nodes = [
    ('入院', 'admission → history_taking → physical_exam → ddx → medication_reconciliation'),
    ('分诊', 'triage → doctor_confirm(人工卡点) → batch_scoring'),
    ('住院', 'monitoring → daily_round → nursing → lab_review'),
    ('出院', 'discharge → handoff → doctor_review(人工卡点) → patient_confirm'),
]
for i, (phase, detail) in enumerate(nodes):
    rrect(s, Inches(0.5), Inches(1.7)+i*Inches(1.2), Inches(5.5), Inches(1.0), HIGHLIGHT, BORDER)
    txt(s, Inches(0.7), Inches(1.8)+i*Inches(1.2), Inches(1.2), Inches(0.4), phase, 14, True, PRIMARY)
    txt(s, Inches(2.0), Inches(1.8)+i*Inches(1.2), Inches(3.8), Inches(0.7), detail, 9, False, GRAY)

# 右侧6道护栏
guards = ['Pydantic输出Schema校验', '幻觉检测(score<0.6→降级)', '模板回退(fallback_to_template)', 'LLM只读原则(不直写DB)', 'source_type溯源(rule/rag/llm/manual)', 'CircuitBreaker三档渐进降级']
section_card(s, Inches(6.5), Inches(1.7), Inches(6.2), Inches(2.4), '6道安全护栏', guards, ACCENT)

# 底部技术细节
tech_details = [
    ('5 Agent体系', '任务规划 → 病历一致性 → 指南合规 → 自检 → 审核编排'),
    ('LLM路由', 'DeepSeek V4-Pro(主力) → Ollama(容灾回退) → 规则引擎(基线)'),
    ('审核卡点', 'doctor_confirm(入院) / med_confirm(调药) / discharge_sign(出院)'),
    ('事件驱动', 'push-based泛型事件循环，planTurn双分支(GEN/RESUME)'),
]
for i, (t, d) in enumerate(tech_details):
    y = Inches(4.4) + i * Inches(0.7)
    txt(s, Inches(6.7), y, Inches(2.0), Inches(0.3), t, 12, True, PRIMARY)
    txt(s, Inches(6.7), y+Inches(0.3), Inches(5.8), Inches(0.35), d, 10, False, GRAY)

# ═══════════════════════════════════════════════════════
# P7: 医生端详解
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '05  医生端功能详解', '6大核心能力 · 85+功能按钮 · 23个临床组件')

# 6步流程横向
steps = [('审核队列', '入院诊断/用药/出院\n三卡点待审核'), ('AI审核', 'DDx编辑/调药决策\n交接事项逐项核对'), ('查房管理', 'SOAP生成/编辑\n体征趋势/证据溯源'), ('医嘱协同', '用药/检查/MDT/\n宣教/随访全流程'), ('照护管理', '生命周期状态流转\n版本化写入+409保护'), ('出院流程', '6步:条件→审核→\n发起→签字→交接→回授')]
for i, (title, desc) in enumerate(steps):
    flow_step(s, Inches(0.3)+i*Inches(2.2), Inches(1.7), Inches(2.0), Inches(1.35), i+1, title, desc, [PRIMARY,PRIMARY_LT,ACCENT,WARNING,SUCCESS,PURPLE][i])

# 核心特性
doc_features_l = [
    ('📋 DiffPanel', '三模式审核(入院/用药/出院) · payload完整暴露 · 编辑优先 · 409冲突保留草稿'),
    ('💊 照护管理', '5类操作(医嘱/检查/MDT/宣教/随访) · 生命周期流转 · canSubmitCareAction校验'),
    ('🔬 监测录入', '体征(血压/心率/SpO2/体温) + 检验(项目/结果/单位) · 趋势可视化 · 异常高亮'),
    ('📝 入院采集', '病史(SOAP/OLDCARTS) + 体格检查(Bates指南10系统) · DDx LLM辅助鉴别'),
    ('🤖 临床助手', '5种模式(查房/护理/用药/患教/中西医) · SSE流式 · 操作草稿审批'),
    ('📄 出院小结', 'AI生成结构化摘要 · PDF导出 · audit_export留痕 · 交接签字'),
]
for i, (title, desc) in enumerate(doc_features_l):
    y = Inches(3.3) + i * Inches(0.65)
    txt(s, Inches(0.6), y, Inches(2.5), Inches(0.25), title, 12, True, PRIMARY)
    txt(s, Inches(3.2), y, Inches(9.5), Inches(0.55), desc, 10, False, GRAY)

# ═══════════════════════════════════════════════════════
# P8: 护士端 + 管理端
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '06-07  护士端 & 管理端详解')

# 护士端 左侧
section_card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(5.3), '护理看板 — 6大工作区',
    ['班次总览: AI优先级排序 + 体征趋势 + 交班摘要',
     '护理任务: 按患者展示待办 · 一键完成 · 录护理',
     '在院患者: 搜索/筛选/风险分层 · 床旁风险提示',
     '逾期监测: 严重逾期红色高亮 · 补录护理入口',
     '交班报告: 重点关注/今日出院/病情稳定三组',
     '制度执行: 4步闭环(标准→任务→升级→留痕)',
     '床旁助手: 护理助手 · 证据图谱 · 临床摘要',
     '护理KPI: 完成率/逾期率/按类型统计'], PRIMARY)

# 管理端 右侧
section_card(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.3), '管理控制台 — 两种角色视图',
    ['【科主任】病区总览: 工作量/趋势/告警/查房',
     '【科主任】知识库治理: 16层385条 · 重建索引',
     '【科主任】病种模板: 14病种 · JSON配置化',
     '【科主任】组织架构: 多科室人员管理',
     '【科主任】系统运维: reindex/seed/clear',
     '【护士长】护理质量: KPI · 完成率 · 交接总览',
     '【护士长】制度执行: 4步闭环 · 确认留痕',
     '【护士长】证据图谱: Neo4j诊断→证据→指南'], ACCENT)

# ═══════════════════════════════════════════════════════
# P9: FHIR 标准化
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '08  FHIR R4 标准化与数据互通', '8类FHIR资源 · Patient Compartment · PII脱敏 · 访问审计')

# FHIR资源卡片
fhir_resources = [
    ('Patient', '患者基本信息/脱敏', PRIMARY),
    ('Encounter', '就诊记录/出入院', PRIMARY_LT),
    ('Condition', '诊断/DDx/ICD-10', ACCENT),
    ('Observation', '体征/检验/LOINC', WARNING),
    ('MedicationRequest', '用药申请/医嘱', SUCCESS),
    ('CarePlan', '照护计划(出院+慢病)', PURPLE),
    ('Consent', '知情同意/授权', DANGER),
    ('AuditEvent', '审计事件/C/R/U/D', GRAY),
]
for i, (name, desc, color) in enumerate(fhir_resources):
    col = i % 4; row = i // 4
    x = Inches(0.5) + col * Inches(3.15)
    y = Inches(1.7) + row * Inches(1.55)
    mini_card(s, x, y, Inches(2.85), Inches(1.3), '', name, desc, WHITE, color)

# 底部特色
fhir_features = [
    ('脱敏输出', 'Name→首字+** · Identifier→TOKEN-后4位'),
    ('访问审计', '每次读/写自动INSERT audit_events · actor/time追踪'),
    ('同步机制', 'inpatient-ward→BackgroundTasks→fhir-adapter'),
    ('出院对接', 'POST /fhir/Observation/Condition/AuditEvent/MedicationRequest'),
]
for i, (t, d) in enumerate(fhir_features):
    x = Inches(0.5) + i * Inches(3.15)
    txt(s, x, Inches(4.8), Inches(2.85), Inches(0.25), t, 12, True, PRIMARY)
    txt(s, x, Inches(5.1), Inches(2.85), Inches(0.5), d, 10, False, GRAY)

# ═══════════════════════════════════════════════════════
# P10: 临床助手
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '09  临床智能助手体系', '5种助手模式 · 16层RAG知识索引 · SSE流式响应')

# 5助手
assistants = [
    ('👨‍⚕️ 查房助手', '主治/住院 · DDx讨论 · DB Agent查病史 · 7层RAG', PRIMARY),
    ('👩‍⚕️ 护理助手', '护士长/护士 · 护理操作 · 交接班 · 清单生成 · 7层RAG', PRIMARY_LT),
    ('💊 用药助手', '医生/药师 · 药物相互作用 · 剂量调整 · OpenFDA实时', ACCENT),
    ('🏠 患教助手', '患者/家属 · 出院须知 · 通俗化改写 · 康复指导', WARNING),
    ('🌿 中西医协同', '医生 · 双视角评估 · 经方速查 · 出院调养 · 6层RAG', PURPLE),
]
for i, (title, desc, color) in enumerate(assistants):
    rrect(s, Inches(0.4), Inches(1.7)+i*Inches(0.9), Inches(6.5), Inches(0.75), LIGHT_BG if i%2==0 else WHITE, BORDER)
    txt(s, Inches(0.6), Inches(1.8)+i*Inches(0.9), Inches(3.5), Inches(0.25), title, 13, True, color)
    txt(s, Inches(4.2), Inches(1.8)+i*Inches(0.9), Inches(2.6), Inches(0.5), desc, 9, False, GRAY)

# 右侧技术架构
section_card(s, Inches(7.4), Inches(1.7), Inches(5.3), Inches(2.8), '技术实现',
    ['POST /assistant/chat · /stream SSE',
     '意图识别 → RAG检索 → 上下文组装 → LLM推理',
     'Redis会话管理(生产) / 内存dict(开发)',
     '30分钟无活动自动清除',
     'ChatBox三模式: 浮动/侧栏/内嵌',
     'SSE首token≤3s 契约'], PRIMARY)

section_card(s, Inches(7.4), Inches(4.8), Inches(5.3), Inches(2.2), '操作草稿闭环',
    ['助手输出→转为操作草稿→医生编辑',
     '审批/驳回→同步到患者状态',
     'PATCH draft · POST approve/reject'], ACCENT)

# ═══════════════════════════════════════════════════════
# P11: 中西医协同 + 知识库
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '10-11  中西医协同 & 知识库体系')

# 左半：中西医
section_card(s, Inches(0.4), Inches(1.7), Inches(6.2), Inches(5.3), '中西医结合智能协同',
    ['L15 中医核心层: 六经辨证公式(8规则) + 经方速查(129+)',
     '药物性味归经(345种) + 倪氏六健康标准',
     '脉诊舌诊速查(16条) + 药食同源禁忌(30条)',
     '三层协同: 入院(体质辨识) → 住院(中西药交互) → 出院(中医康复)',
     '三省策略: 诊断公式 + 药物协同 + 体质出院评估',
     '3批次实施: S1知识入库 → S2药物交互 → S3体质评估',
     '严格边界: 不蒸馏345种本草/849医案，不接入非医学内容'], PURPLE)

# 右半：知识库
section_card(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.3), 'RAG知识库 & 证据图谱',
    ['16层 / 385条结构化临床知识',
     'Milvus 2.4 向量检索 (dim=768, IP度量)',
     '知识生命周期: 7种状态枚举 + 状态机流转',
     '语义检索预览 + 分层索引管理',
     'Neo4j 证据图谱: 诊断→证据→指南路径',
     'RAG→LLM管线: prompt拼接+字段级验证',
     '4批次知识治理: 导入→校验→发布→反向阻断'], PRIMARY)

# ═══════════════════════════════════════════════════════
# P12: 安全审计
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '12  安全与审计体系', '三模式认证 · 角色中间件 · 全链路追踪 · 审计不可覆盖')

# 鉴权
auth_items = [
    ('🔐 三模式认证', 'header(开发) / jwt(联调) / oidc(生产Keycloak)\nAccess Token 15min · Refresh Token 8h'),
    ('👥 角色中间件', 'doctor/nurse双角色 · 端点级RBAC\n医生专属路由(审核/出院/command)\n患者级访问控制(patient_access)'),
    ('🛡️ 数据安全', 'PII脱敏(name/identifier)\nFHIR访问审计(C/R/U/D四操作)\n请求ID全链路追踪(X-Request-ID)'),
    ('📝 审计体系', 'audit_events 不可变记录\nbefore_state/after_state 变更追踪\nactor + timestamp 逐条记录'),
]
for i, (title, desc) in enumerate(auth_items):
    y = Inches(1.8) + i * Inches(1.3)
    rrect(s, Inches(0.5), y, Inches(6.5), Inches(1.1), HIGHLIGHT if i%2==0 else WHITE, BORDER)
    txt(s, Inches(0.7), y+Inches(0.1), Inches(6.0), Inches(0.3), title, 14, True, PRIMARY)
    txt(s, Inches(0.7), y+Inches(0.45), Inches(6.0), Inches(0.55), desc, 10, False, GRAY)

# 右侧安全功能
section_card(s, Inches(7.5), Inches(1.8), Inches(5.3), Inches(5.2), '安全技术栈',
    ['Idempotency-Key 幂等中间件',
     'expected_version CAS乐观锁',
     '409冲突后保留用户草稿',
     '不自动回放临床写入',
     'LLM只读原则 · Harness护栏',
     '请求级X-Request-ID追踪',
     '结构化JSON日志 · HTTP指标',
     'SKIP_BRIDGE外网阻断开关',
     '生产环境Fixture禁用',
     'DeepSeek API 401→Ollama回退'], ACCENT)

# ═══════════════════════════════════════════════════════
# P13: 质量保障
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '13  质量保障与测试体系')

# 测试金字塔
levels = [
    ('🧪 E2E (17条)', '3条核心链路 / Playwright自动化 / 全覆盖', 2.2),
    ('🔬 前端单元 (128条)', '46测试文件 / React组件+Services+Hooks', 2.9),
    ('⚙️ 后端Agent (48条)', '19节点独立测试 / Mock LLM确定性', 3.6),
    ('📋 接口契约', 'TS编译零错误 / OpenAPI自检 / shared contracts', 4.3),
]
for icon, desc, y_in in levels:
    y = Inches(y_in)
    rrect(s, Inches(0.8), y, Inches(5.5), Inches(0.6), WHITE, BORDER)
    txt(s, Inches(1.0), y+Inches(0.08), Inches(5.0), Inches(0.2), f'{icon}', 14, True, PRIMARY)
    txt(s, Inches(1.0), y+Inches(0.3), Inches(5.0), Inches(0.25), desc, 9, False, GRAY)

# 右侧
eng_items = [
    'TypeScript 全栈类型安全',
    'React Query 缓存 + 乐观更新',
    'Vitest + Playwright 双测试框架',
    'Python 3.13 + asyncio 异步',
    'SQLAlchemy 2.0 ORM',
    'Pydantic v2 校验',
    'Vite code splitting + manualChunks',
    'React.memo + useCallback',
    'PanelErrorBoundary 面板级容错',
    'Git版本化Prompt · A/B测试',
    'CI/CD: ruff+mypy+pytest',
    'Docker Compose 5容器',
]
section_card(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.3), '工程实践', eng_items, SUCCESS)

# ═══════════════════════════════════════════════════════
# P14: 项目价值
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, PRIMARY)
rect(s, Inches(0), Inches(0), Inches(0.12), SLIDE_H, PRIMARY_LT)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.12), PRIMARY_LT)
txt(s, Inches(1.2), Inches(0.6), Inches(11), Inches(0.6), '14  项目价值与预期效果', 32, True, WHITE)

vals = [
    ('📉', '降低再入院率', '出院后7天真空期标准化交接\n随访闭环覆盖 · 预期降低20-30%'),
    ('⏱️', '提升临床效率', 'AI查房摘要+出院小结自动生成\n医生审核时间缩短40%'),
    ('🔄', '打通数据孤岛', 'FHIR R4标准 · 出院数据一键同步\n慢病管理数据连续不中断'),
    ('🛡️', '降低医疗风险', 'NEWS2告警+用药冲突检测\nAI辅助DDx+全链路审计'),
]
for i, (icon, title, desc) in enumerate(vals):
    x = Inches(0.5) + i * Inches(3.2)
    txt(s, x, Inches(1.8), Inches(2.8), Inches(0.5), icon, 36, False, WHITE, PP_ALIGN.CENTER)
    txt(s, x, Inches(2.5), Inches(2.8), Inches(0.4), title, 18, True, WHITE, PP_ALIGN.CENTER)
    rect(s, x+Inches(0.5), Inches(3.0), Inches(1.8), Inches(0.03), ACCENT)
    txt(s, x, Inches(3.3), Inches(2.8), Inches(1.5), desc, 11, False, RGBColor(0xCC,0xE0,0xE5), PP_ALIGN.CENTER)

# 底部指标
metrics = [('130+', '后端测试'), ('128', '前端测试'), ('17', 'E2E链路'), ('85+', '医生端按钮'), ('19', 'AI节点'), ('163+', 'API端点')]
for i, (n, l) in enumerate(metrics):
    x = Inches(0.8) + i * Inches(2.1)
    txt(s, x, Inches(5.2), Inches(1.8), Inches(0.5), n, 30, True, ACCENT, PP_ALIGN.CENTER)
    txt(s, x, Inches(5.8), Inches(1.8), Inches(0.3), l, 11, False, RGBColor(0xBB,0xD5,0xDC), PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# P15: 路线图
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '15  实施路线图', '三阶段交付 · 从试点到规模化')

phases = [
    ('✅ 已完成 v0.3', SUCCESS, [
        '医生工作台(85+功能按钮)', '护理看板(6大工作区)', '管理控制台(双角色视图)',
        '19节点AI Agent编排', 'FHIR Adapter接入主链路', '16层385条RAG知识',
        'Neo4j证据图谱', 'E2E自动化测试(17条)', '3个病种模板(心衰/高血压/糖尿病)',
    ]),
    ('🔄 进行中 v0.4', ACCENT, [
        'Docker Compose 一键部署', 'WebSocket实时推送', 'API限流与安全加固',
        '审计日志MySQL持久化', 'Alembic数据库迁移', '14病种模板全覆盖',
        '中医知识库L15层', 'carehandoff对接社区', 'Node.js→Python全线迁移',
    ]),
    ('📋 规划中 v1.0', PRIMARY, [
        '患者端(微信小程序)', '多院区组织架构扩展', '短信/微信随访提醒',
        'BI质控数据看板', 'CDSS临床决策认证', 'Keycloak OIDC生产部署',
        'K8s+Ingress部署', '金标准临床评测集', '区域转诊协同网络',
    ]),
]
for i, (stage, color, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(4.3)
    rrect(s, x, Inches(1.8), Inches(3.9), Inches(0.5), color)
    txt(s, x, Inches(1.88), Inches(3.9), Inches(0.35), stage, 15, True, WHITE, PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        txt(s, x+Inches(0.15), Inches(2.5)+j*Inches(0.42), Inches(3.7), Inches(0.38), f'• {item}', 10, False, BLACK)

# ═══════════════════════════════════════════════════════
# P16: 技术债务与风险
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '16  当前架构评分与待办事项', 'R3架构审查: 8.0/10 · 130/130测试通过 · 19节点全链路覆盖')

dims = [
    ('模块边界', '7.0→8.5', '+1.5', SUCCESS),
    ('代码复杂度', '6.5→8.0', '+1.5', SUCCESS),
    ('依赖关系', '7.0→8.0', '+1.0', SUCCESS),
    ('LLM集成', '5.5→8.5', '+3.0', SUCCESS),
    ('新增模块质量', '8.0', '—', WARNING),
]
for i, (name, score, delta, color) in enumerate(dims):
    x = Inches(0.5) + i * Inches(2.55)
    kpi_card(s, x, Inches(1.7), Inches(2.25), Inches(1.1), score, f'{name} ({delta})', color)

debts = [
    ('P0 本周', DANGER, ['鉴权升级(header→JWT+Keycloak)', '3处跨层依赖修复(agent→routes)', '路由层绕过领域层解耦']),
    ('P1 本月', WARNING, ['state_store持久化(SqliteSaver)', 'LLM节点优化(8→4必要)', 'API补齐(批量端点/FHIR全写)']),
    ('P2 下季度', PRIMARY, ['模板深化(14→22病种)', '可观测性(Prometheus+Grafana)', 'Docker化部署自动化']),
]
for i, (priority, color, items) in enumerate(debts):
    y = Inches(3.1) + i * Inches(1.4)
    rrect(s, Inches(0.5), y, Inches(12.3), Inches(1.25), WHITE, BORDER)
    txt(s, Inches(0.7), y+Inches(0.08), Inches(1.5), Inches(0.3), priority, 13, True, color)
    for j, item in enumerate(items):
        txt(s, Inches(2.4)+j*Inches(3.5), y+Inches(0.08), Inches(3.2), Inches(0.25), f'• {item}', 10, False, GRAY)

# ═══════════════════════════════════════════════════════
# P17: 竞品对比
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, WHITE)
title_bar(s, '17  技术亮点一览', '六大核心能力 vs 行业现状')

highlights = [
    ('🏗️ 微服务架构', '5服务独立部署 · HTTP REST通信\n统一UnifiedResponse · request_id透传', '传统单体HIS耦合严重\n接口格式不统一'),
    ('🤖 Agent编排', '19节点StateGraph · 双模式\nDeepSeek V4 + Ollama容灾', '规则引擎为主\n缺乏AI推理能力'),
    ('📊 证据溯源', '每条临床建议附带RAG引用\nNeo4j图谱可视化路径', '黑盒决策\n无法追溯推理过程'),
    ('🔒 并发安全', 'CAS乐观锁+事务状态机\n409冲突保留用户草稿', '后写覆盖\n多人并发写入丢失'),
    ('✅ 测试体系', '130后端+128前端+17E2E\nTS零错误 · 全角色覆盖', '测试不足\n生产问题频发'),
    ('🏥 FHIR标准', '8类资源+Patient Compartment\n脱敏+审计+同步', '私有格式\n系统间无法互通'),
]
for i, (icon, zhenhu, others) in enumerate(highlights):
    row = i // 3; col = i % 3
    x = Inches(0.4) + col * Inches(4.2)
    y = Inches(1.7) + row * Inches(2.7)
    rrect(s, x, y, Inches(3.9), Inches(2.45), WHITE, BORDER)
    txt(s, x+Inches(0.15), y+Inches(0.1), Inches(3.6), Inches(0.3), icon, 22)
    txt(s, x+Inches(0.15), y+Inches(0.45), Inches(3.6), Inches(0.55), zhenhu, 11, False, BLACK)
    txt(s, x+Inches(0.15), y+Inches(1.3), Inches(3.6), Inches(0.25), 'vs', 9, True, ACCENT)
    txt(s, x+Inches(0.15), y+Inches(1.6), Inches(3.6), Inches(0.6), others, 10, False, GRAY)

# ═══════════════════════════════════════════════════════
# P18: 感谢页
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, PRIMARY)
rect(s, Inches(0), Inches(0), Inches(0.12), SLIDE_H, PRIMARY_LT)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.12), PRIMARY_LT)
txt(s, Inches(1.2), Inches(2.2), Inches(11), Inches(1.0), '感谢聆听', 60, True, WHITE, PP_ALIGN.CENTER)
rect(s, Inches(5.0), Inches(3.5), Inches(3.5), Inches(0.04), ACCENT)
txt(s, Inches(1.2), Inches(4.0), Inches(11), Inches(0.6), '臻护 — 让每一次出院交接都有据可循', 22, False, RGBColor(0xCC,0xE0,0xE5), PP_ALIGN.CENTER)
txt(s, Inches(1.2), Inches(5.5), Inches(11), Inches(0.4), '心内科试点  ·  臻护团队  ·  2026年7月', 14, False, GRAY, PP_ALIGN.CENTER)
txt(s, Inches(1.2), Inches(6.2), Inches(11), Inches(0.4), '基于46份设计文档 · 163+ API · 19节点AI Agent · FHIR R4 · 全角色覆盖\n欢迎提问与交流', 11, False, RGBColor(0x99,0xBB,0xC5), PP_ALIGN.CENTER)

# ── 保存 ──
prs.save(OUTPUT)
print(f'PPT v2.0 已生成 ({prs.slides.__len__()} 页): {OUTPUT}')
